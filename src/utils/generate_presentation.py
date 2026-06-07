import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color Palette (Sober & Flat Sapphire Theme)
COLOR_SAPPHIRE = RGBColor(59, 130, 246)    # #3B82F6 - Primary Sapphire Blue
COLOR_ROYAL = RGBColor(29, 78, 216)        # #1D4ED8 - Darker Royal Blue Accent
COLOR_CHARCOAL = RGBColor(31, 41, 55)      # #1F2937 - Body Text Dark
COLOR_MUTED = RGBColor(107, 114, 128)      # #6B7280 - Captions & Footers Gray
COLOR_LIGHT_BG = RGBColor(243, 244, 246)   # #F3F4F6 - Light Panel Gray
COLOR_WHITE = RGBColor(255, 255, 255)      # #FFFFFF - Slide Background

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

def add_header(slide, title_text):
    """Adds a standard academic header to the slide."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_SAPPHIRE
    
    # Add a thin flat divider line below header
    shape = slide.shapes.add_shape(
        1, # Rectangle
        Inches(0.5), Inches(1.15), Inches(12.333), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ROYAL
    shape.line.color.rgb = COLOR_ROYAL
    shape.line.width = Pt(0)

def add_footer(slide, current_slide, total_slides=10):
    """Adds a standard professional footer to the slide."""
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.3))
    tf = footer_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"BUT2 Internship Defense | Milan Loi | Slide {current_slide}/{total_slides}"
    p.font.name = FONT_BODY
    p.font.size = Pt(9.5)
    p.font.color.rgb = COLOR_MUTED

def set_slide_background(slide):
    """Sets a clean white background for the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

def create_presentation():
    # Read modularity Q score, node and edge count dynamically if available
    modularity_q = "0.9539" # Baseline actual value
    node_count = "20641"
    edge_count = "18087"
    metrics_path = "deliverables/week_5/echo_chamber_metrics.txt"
    if os.path.exists(metrics_path):
        try:
            with open(metrics_path, "r") as f:
                for line in f:
                    if "Modularity Score (Q):" in line:
                        modularity_q = line.split(":")[-1].strip()
                    elif "Total Core Nodes:" in line:
                        node_count = line.split(":")[-1].strip()
                    elif "Total Core Edges:" in line:
                        edge_count = line.split(":")[-1].strip()
        except Exception:
            pass

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    total_slides = 10
    
    # ==========================================
    # SLIDE 1: TITLE SLIDE (Cover Page)
    # ==========================================
    slide_1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_1)
    
    # Title Box (centered)
    title_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CROSS-PLATFORM STUDY AND VISUALIZATION OF DISCURSIVE POLARIZATION"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_SAPPHIRE
    
    p2 = tf.add_paragraph()
    p2.text = "Modeling In-Group vs Out-Group Tribal Opposition Mechanisms on Social Media"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ROYAL
    p2.space_before = Pt(8)

    # Confidential notice
    p_conf = tf.add_paragraph()
    p_conf.text = "MENTION : CONFIDENTIEL"
    p_conf.alignment = PP_ALIGN.CENTER
    p_conf.font.name = FONT_TITLE
    p_conf.font.size = Pt(14)
    p_conf.font.bold = True
    p_conf.font.color.rgb = COLOR_ROYAL
    p_conf.space_before = Pt(8)

    # Info Box
    info_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.333), Inches(2.5))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    p_info = tf_info.paragraphs[0]
    p_info.text = "Student: Milan LOI | BUT Informatique (Parcours Réalisation d'applications)\nUniversity: Université de Limoges - IUT de Limoges, Département Informatique\nHost Institution: Robert Gordon University (RGU), School of Computing Science and Digital Media"
    p_info.alignment = PP_ALIGN.CENTER
    p_info.font.name = FONT_BODY
    p_info.font.size = Pt(11)
    p_info.font.color.rgb = COLOR_CHARCOAL
    p_info.space_after = Pt(6)
    
    p_tut = tf_info.add_paragraph()
    p_tut.text = "Host Supervisor: Dr. Shahana Bano (RGU) | Academic Advisor: M. Jean-Marc Garcia (IUT de Limoges)\nInternship Period: April 6, 2026 – June 12, 2026 (10 Weeks)"
    p_tut.alignment = PP_ALIGN.CENTER
    p_tut.font.name = FONT_BODY
    p_tut.font.size = Pt(10.5)
    p_tut.font.italic = True
    p_tut.font.color.rgb = COLOR_MUTED
    
    add_footer(slide_1, 1, total_slides)
    
    # ==========================================
    # SLIDE 2: DEFENSE PLAN (AGENDA)
    # ==========================================
    slide_2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_2)
    add_header(slide_2, "Defense Agenda")
    add_footer(slide_2, 2, total_slides)
    
    left_box = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.9), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_lh = tf_left.paragraphs[0]
    p_lh.text = "Part I: Context & Preparatory Steps"
    p_lh.font.name = FONT_TITLE
    p_lh.font.size = Pt(18)
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_ROYAL
    p_lh.space_after = Pt(10)
    
    l_points = [
        ("1. Introduction & Research Question", "The societal threat of echo chambers and verbal hostility."),
        ("2. Host Institution Profile", "NLP and Social Data Science Research Lab at Robert Gordon University."),
        ("3. Comparative Solution Analysis", "Justifying the transition from lexical baselines to deep context models."),
        ("4. System Architecture & Ingestion", "Building the processing pipeline, cleaning, and language filtering.")
    ]
    for title, desc in l_points:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(11.5)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(6)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10.5)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL
        
    # Right panel
    panel_agenda = slide_2.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_agenda.fill.solid()
    panel_agenda.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_agenda.line.color.rgb = COLOR_ROYAL
    panel_agenda.line.width = Pt(1.5)
    
    right_box = slide_2.shapes.add_textbox(Inches(6.85), Inches(1.6), Inches(5.8), Inches(4.8))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_rh = tf_right.paragraphs[0]
    p_rh.text = "Part II: Execution & Reflexive Review"
    p_rh.font.name = FONT_TITLE
    p_rh.font.size = Pt(18)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_ROYAL
    p_rh.space_after = Pt(10)
    
    r_points = [
        ("5. Modality & Topological Graph Models", "CLIP zero-shot alignment and Louvain modularity clustering."),
        ("6. Technical Challenge & Adaptive Resolution", "Resolving Reddit API PRAW authentication blocks via Hugging Face."),
        ("7. Empirical Findings & Visual Dashboard", "Verifying polarization score trends and WebGL dashboard layers."),
        ("8. BUT Skills & Critical Reflection", "Self-evaluation of the four target BUT competencies and SWOT matrix."),
        ("9. AI Responsibility Disclosure", "Detailed logging of prompts and model utilization parameters (<15%).")
    ]
    for title, desc in r_points:
        p_pt = tf_right.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(11.5)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(4)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10.5)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL
        
    # ==========================================
    # SLIDE 3: INTRODUCTION & RESEARCH QUESTION
    # ==========================================
    slide_3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_3)
    add_header(slide_3, "1. Introduction & Research Problématique")
    add_footer(slide_3, 3, total_slides)
    
    left_box = slide_3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_c = tf_left.paragraphs[0]
    p_c.text = "Societal Context & Rising Polarisation"
    p_c.font.name = FONT_TITLE
    p_c.font.size = Pt(18)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_ROYAL
    p_c.space_after = Pt(12)
    
    c_points = [
        ("Tribal Segregation", "Algorithmic amplification creates rigid 'Us vs Them' dynamics on web platforms."),
        ("Echo Chambers Presence", "Empirical data shows 65%+ of discussions occur in isolated clusters, with less than 5% exposure to opposing views."),
        ("Need for Quantification", "Manual auditing fails to capture large scale behaviors, making automated cross-platform pipelines necessary.")
    ]
    for title, desc in c_points:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(12)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL

    # Right panel - Problématique
    panel_prob = slide_3.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_prob.fill.solid()
    panel_prob.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_prob.line.color.rgb = COLOR_ROYAL
    panel_prob.line.width = Pt(1.5)
    
    right_box = slide_3.shapes.add_textbox(Inches(6.85), Inches(2.2), Inches(5.8), Inches(4.0))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_prh = tf_right.paragraphs[0]
    p_prh.text = "FORMAL RESEARCH QUESTION (PROBLÉMATIQUE)"
    p_prh.alignment = PP_ALIGN.CENTER
    p_prh.font.name = FONT_TITLE
    p_prh.font.size = Pt(16)
    p_prh.font.bold = True
    p_prh.font.color.rgb = COLOR_ROYAL
    p_prh.space_after = Pt(20)
    
    p_pq = tf_right.add_paragraph()
    p_pq.text = "How can we quantify, model topologically, and visualize in an integrated platform the linguistic polarization markers (We vs Them tribal dynamics) and echo chamber boundaries across heterogeneous social media sources?"
    p_pq.alignment = PP_ALIGN.CENTER
    p_pq.font.name = FONT_BODY
    p_pq.font.size = Pt(14)
    p_pq.font.italic = True
    p_pq.font.bold = True
    p_pq.font.color.rgb = COLOR_CHARCOAL
    
    # ==========================================
    # SLIDE 4: HOST INSTITUTION
    # ==========================================
    slide_4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_4)
    add_header(slide_4, "2. Host Institution: Robert Gordon University")
    add_footer(slide_4, 4, total_slides)
    
    left_box = slide_4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_inst = tf_left.paragraphs[0]
    p_inst.text = "School of Computing & Digital Media"
    p_inst.font.name = FONT_TITLE
    p_inst.font.size = Pt(18)
    p_inst.font.bold = True
    p_inst.font.color.rgb = COLOR_ROYAL
    p_inst.space_after = Pt(12)
    
    inst_points = [
        ("Location & Reputation", "Located in Aberdeen, Scotland (Garthdee campus). Ranked top in Scotland for graduate employability."),
        ("Applied Research focus", "Strong partnerships with international academic networks and tech industries."),
        ("Scientific Standards", "RGU enforces strict validation frameworks, reproducible data setups, and ethical AI utilization.")
    ]
    for title, desc in inst_points:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(12)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL

    # Right panel - Team and Advisors
    panel_team = slide_4.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_team.fill.solid()
    panel_team.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_team.line.color.rgb = COLOR_ROYAL
    panel_team.line.width = Pt(1.5)
    
    right_box = slide_4.shapes.add_textbox(Inches(6.85), Inches(1.7), Inches(5.8), Inches(4.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_team_h = tf_right.paragraphs[0]
    p_team_h.text = "NLP & Social Data Science Team"
    p_team_h.font.name = FONT_TITLE
    p_team_h.font.size = Pt(18)
    p_team_h.font.bold = True
    p_team_h.font.color.rgb = COLOR_ROYAL
    p_team_h.space_after = Pt(12)
    
    team_points = [
        ("Host Supervisor", "Dr. Shahana Bano: Researcher in Natural Language Processing, machine learning, and cybersecurity."),
        ("Infrastructure", "Access to GPU-accelerated computing nodes for local deep learning inference (transformers, CLIP embeddings)."),
        ("Collaboration", "Participated in weekly research sprints with post-docs and PhD candidates to refine data paradigms.")
    ]
    for title, desc in team_points:
        p_pt = tf_right.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(11.5)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10.5)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL
        
    # ==========================================
    # SLIDE 5: TECHNICAL SOLUTIONS EVALUATED
    # ==========================================
    slide_5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_5)
    add_header(slide_5, "3. Comparative Solution Analysis")
    add_footer(slide_5, 5, total_slides)
    
    left_box = slide_5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_comp = tf_left.paragraphs[0]
    p_comp.text = "Methodological Comparisons & Benchmarks"
    p_comp.font.name = FONT_TITLE
    p_comp.font.size = Pt(18)
    p_comp.font.bold = True
    p_comp.font.color.rgb = COLOR_ROYAL
    p_comp.space_after = Pt(12)
    
    comp_points = [
        ("Data Ingestion", "Official APIs (PRAW/yt-dlp) are flexible but restricted. Pre-compiled Hugging Face datasets ensure high volume and reproducibility."),
        ("Sentiment / Toxicity Assessment", "Traditional keyword list counting (VADER) returned near-zero correlation coefficients. Google's deep Perspective API captures context, grammar, and intent."),
        ("Visual Dashboard Engine", "Classic JS/D3.js frontend requires a separate REST backend. Streamlit provides direct Python integration, fast data caching, and native Plotly wrappers.")
    ]
    for title, desc in comp_points:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(12)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL

    # Right side table - Comparison of Sentiment Models
    panel_tbl = slide_5.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_tbl.fill.solid()
    panel_tbl.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_tbl.line.color.rgb = COLOR_ROYAL
    panel_tbl.line.width = Pt(1.5)
    
    right_box = slide_5.shapes.add_textbox(Inches(6.85), Inches(1.7), Inches(5.8), Inches(4.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_tbl_h = tf_right.paragraphs[0]
    p_tbl_h.text = "Benchmark: Keyword Counting vs Deep NLP"
    p_tbl_h.font.name = FONT_TITLE
    p_tbl_h.font.size = Pt(16)
    p_tbl_h.font.bold = True
    p_tbl_h.font.color.rgb = COLOR_ROYAL
    p_tbl_h.space_after = Pt(15)
    
    # Add small table in presentation
    table_shape = slide_5.shapes.add_table(4, 3, Inches(6.85), Inches(2.5), Inches(5.8), Inches(2.2))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    
    headers = ["Metric / Method", "VADER (Lexical)", "Perspective (Deep)"]
    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ROYAL
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    t_data = [
        ["Spearman Correlation (Reddit)", "ρ = 0.0148 (p=0.22)", "ρ = 0.3842 (p<0.001)"],
        ["Sarcasm/Slang", "Blind", "Context-Aware"],
        ["Execution Speed", "Ultra-Fast (CPU)", "API Throttled / GPU"]
    ]
    for r_idx, row in enumerate(t_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_CHARCOAL
            if c_idx == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.bold = True
            else:
                p.alignment = PP_ALIGN.CENTER
                
    cap_box = slide_5.shapes.add_textbox(Inches(6.85), Inches(5.0), Inches(5.8), Inches(1.0))
    tf_cap = cap_box.text_frame
    tf_cap.word_wrap = True
    p_cap = tf_cap.paragraphs[0]
    p_cap.text = "Result: Lexical word lists failed to achieve statistical significance. Context-aware models (Perspective, CLIP) are required to capture true discursive hostility."
    p_cap.font.name = FONT_BODY
    p_cap.font.size = Pt(10)
    p_cap.font.italic = True
    p_cap.font.color.rgb = COLOR_MUTED
    
    # ==========================================
    # SLIDE 6: SYSTEM ARCHITECTURE & IMPLEMENTATION
    # ==========================================
    slide_6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_6)
    add_header(slide_6, "4. Ingestion & Preprocessing Pipeline")
    add_footer(slide_6, 6, total_slides)
    
    left_box = slide_6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_arch = tf_left.paragraphs[0]
    p_arch.text = "Standardisation & Ingestion Flow"
    p_arch.font.name = FONT_TITLE
    p_arch.font.size = Pt(18)
    p_arch.font.bold = True
    p_arch.font.color.rgb = COLOR_ROYAL
    p_arch.space_after = Pt(12)
    
    arch_points = [
        ("Clean and Deduplicate", "Removes missing fields and duplicate comments induced by sharing networks."),
        ("Language Detection", "Integrates 'langdetect' module to isolate English text ('en'), ensuring downstream model compatibility."),
        ("Time Normalisation", "Warp and map Reddit/YouTube timelines into parallel 2018-2026 timelines for synchronization.")
    ]
    for title, desc in arch_points:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(12)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL

    # Right side panel - Code snippet of Cleaning
    panel_snip = slide_6.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_snip.fill.solid()
    panel_snip.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_snip.line.color.rgb = COLOR_ROYAL
    panel_snip.line.width = Pt(1.5)
    
    right_box = slide_6.shapes.add_textbox(Inches(6.85), Inches(1.7), Inches(5.8), Inches(4.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_snip_h = tf_right.paragraphs[0]
    p_snip_h.text = "Python Clean Pipeline Module"
    p_snip_h.font.name = FONT_TITLE
    p_snip_h.font.size = Pt(16)
    p_snip_h.font.bold = True
    p_snip_h.font.color.rgb = COLOR_ROYAL
    p_snip_h.space_after = Pt(15)
    
    p_code = tf_right.add_paragraph()
    p_code.text = (
        "def clean_dataset(df):\n"
        "    # Remove nulls & short texts\n"
        "    df = df.dropna(subset=['text'])\n"
        "    df = df[df['text'].str.len() > 3]\n\n"
        "    # Deduplicate sémantiquement\n"
        "    df = df.drop_duplicates(subset=['text'])\n\n"
        "    # Filter language\n"
        "    df['lang'] = df['text'].apply(safe_detect_lang)\n"
        "    return df[df['lang'] == 'en']"
    )
    p_code.font.name = "Courier New"
    p_code.font.size = Pt(10.5)
    p_code.font.color.rgb = COLOR_CHARCOAL
    
    # ==========================================
    # SLIDE 7: MODALITY & TOPOLOGICAL GRAPH MODELS
    # ==========================================
    slide_7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_7)
    add_header(slide_7, "5. Modality & Topological Graph Models")
    add_footer(slide_7, 7, total_slides)
    
    left_box = slide_7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_mod = tf_left.paragraphs[0]
    p_mod.text = "Multimodality (CLIP) & Network Topologies"
    p_mod.font.name = FONT_TITLE
    p_mod.font.size = Pt(18)
    p_mod.font.bold = True
    p_mod.font.color.rgb = COLOR_ROYAL
    p_mod.space_after = Pt(12)
    
    mod_points = [
        ("CLIP Zero-Shot Embeddings", "Evaluates text-image similarity using OpenAI's CLIP model. Classifies memes into: in-group, out-group, or neutral."),
        ("Bipartite Network Models", "Constructs bipartite graphs linking users to interaction sources. Prunes degree-1 leaf nodes to keep dense core groups."),
        ("Louvain Community Detection", "Maximizes network modularity (Q) to partition nodes into distinct sub-clusters automatically.")
    ]
    for title, desc in mod_points:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(12)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL

    # Right panel - Math Equations
    panel_math = slide_7.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_math.fill.solid()
    panel_math.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_math.line.color.rgb = COLOR_ROYAL
    panel_math.line.width = Pt(1.5)
    
    right_box = slide_7.shapes.add_textbox(Inches(6.85), Inches(1.7), Inches(5.8), Inches(4.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_math_h = tf_right.paragraphs[0]
    p_math_h.text = "Underlying Mathematical Models"
    p_math_h.font.name = FONT_TITLE
    p_math_h.font.size = Pt(18)
    p_math_h.font.bold = True
    p_math_h.font.color.rgb = COLOR_ROYAL
    p_math_h.space_after = Pt(15)
    
    p_eq1_title = tf_right.add_paragraph()
    p_eq1_title.text = "1. Louvain Modularity Maximisation (Q):"
    p_eq1_title.font.name = FONT_BODY
    p_eq1_title.font.size = Pt(12)
    p_eq1_title.font.bold = True
    p_eq1_title.font.color.rgb = COLOR_ROYAL
    p_eq1_title.space_before = Pt(6)
    
    p_eq1 = tf_right.add_paragraph()
    p_eq1.text = "Q = (1 / 2m) * ∑_ij [ A_ij - (k_i * k_j / 2m) ] * δ(c_i, c_j)"
    p_eq1.font.name = "Consolas"
    p_eq1.font.size = Pt(11)
    p_eq1.font.color.rgb = COLOR_CHARCOAL
    p_eq1.space_before = Pt(4)
    p_eq1.space_after = Pt(12)

    p_eq2_title = tf_right.add_paragraph()
    p_eq2_title.text = "2. Pronoun-Based Polarization Index (PI_c):"
    p_eq2_title.font.name = FONT_BODY
    p_eq2_title.font.size = Pt(12)
    p_eq2_title.font.bold = True
    p_eq2_title.font.color.rgb = COLOR_ROYAL
    
    p_eq2 = tf_right.add_paragraph()
    p_eq2.text = "PI_c = (Frequency(We) / Frequency(Them)) * Mean_Toxicity_c"
    p_eq2.font.name = "Consolas"
    p_eq2.font.size = Pt(11)
    p_eq2.font.color.rgb = COLOR_CHARCOAL
    p_eq2.space_before = Pt(4)
    
    # ==========================================
    # SLIDE 8: TECHNICAL CHALLENGE & ADAPTIVE RESOLUTION
    # ==========================================
    slide_8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_8)
    add_header(slide_8, "6. Technical Challenge & Adaptive Resolution")
    add_footer(slide_8, 8, total_slides)
    
    left_box = slide_8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_ch = tf_left.paragraphs[0]
    p_ch.text = "The Reddit API Block (PRAW)"
    p_ch.font.name = FONT_TITLE
    p_ch.font.size = Pt(18)
    p_ch.font.bold = True
    p_ch.font.color.rgb = COLOR_ROYAL
    p_ch.space_after = Pt(12)
    
    ch_points = [
        ("Authentication Block", "Strict API access limits and OAuth token rejection by Reddit in week 1 threatened the scraping schedule."),
        ("Impact on Timelines", "Risked leaving the project without political Reddit data, invalidating the planned cross-platform benchmark."),
        ("Initial Mitigation", "Attempted rate throttling and user agent rotations, which did not resolve the core access refusal.")
    ]
    for title, desc in ch_points:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(12)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL

    # Right panel - Resolution
    panel_res = slide_8.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_res.fill.solid()
    panel_res.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_res.line.color.rgb = COLOR_ROYAL
    panel_res.line.width = Pt(1.5)
    
    right_box = slide_8.shapes.add_textbox(Inches(6.85), Inches(1.7), Inches(5.8), Inches(4.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_res_h = tf_right.paragraphs[0]
    p_res_h.text = "Adaptive Pivot: Hugging Face Datasets"
    p_res_h.font.name = FONT_TITLE
    p_res_h.font.size = Pt(18)
    p_res_h.font.bold = True
    p_res_h.font.color.rgb = COLOR_ROYAL
    p_res_h.space_after = Pt(12)
    
    res_points = [
        ("Dynamic Integration", "Switched to pre-compiled, peer-reviewed academic datasets hosted on Hugging Face (e.g., mo-mittal/reddit_political_subs)."),
        ("Timeline Alignment", "Filtered and synchronized dataset timestamps to match the 2018-2026 timeline of the YouTube comments."),
        ("Key Learning", "Demonstrated Agile adaptability in the face of API blocks, prioritizing project continuity and statistical volume.")
    ]
    for title, desc in res_points:
        p_pt = tf_right.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(11.5)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(8)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10.5)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL
        
    # ==========================================
    # SLIDE 9: RESULTS & FINDINGS (OLD SLIDE 5)
    # ==========================================
    slide_9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_9)
    add_header(slide_9, "7. Empirical Results & Visual Findings")
    add_footer(slide_9, 9, total_slides)
    
    # Left text description
    left_box = slide_9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_res_title = tf_left.paragraphs[0]
    p_res_title.text = "Key Empirical Discoveries"
    p_res_title.font.name = FONT_TITLE
    p_res_title.font.size = Pt(18)
    p_res_title.font.bold = True
    p_res_title.font.color.rgb = COLOR_ROYAL
    p_res_title.space_after = Pt(10)
    
    bullets = [
        "Lexical Baseline Failure: Spearman rank correlation between negative word counts and engagement was near zero (0.0148 on Reddit, 0.0050 on YouTube).",
        f"Echo Chamber Modularity: The merged cross-platform network yielded modularity Q = {modularity_q} across {node_count} core nodes and {edge_count} edges, proving severe community partitioning.",
        "Tribalism & Toxicity: Active partisan clusters (Alt-Right, Far-Left) exhibit elevated We/Them ratios (>= 1.0) correlating with higher Perspective toxicity.",
        "Cross-Platform Contrast: Reddit shows distinct in-group/out-group language spikes, whereas YouTube comment threads show higher baseline toxicity."
    ]
    
    for bullet in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = "• " + bullet
        p_b.font.name = FONT_BODY
        p_b.font.size = Pt(11.2)
        p_b.font.color.rgb = COLOR_CHARCOAL
        p_b.space_before = Pt(5)
        if f"Q = {modularity_q}" in bullet or "Spearman rank correlation" in bullet or "Cross-Platform Contrast" in bullet:
            p_b.font.bold = True
            
    # Right Image 1: Polarization Heatmap (Top)
    img_path1 = "deliverables/week_6/cluster_polarization_heatmap.png"
    if os.path.exists(img_path1):
        slide_9.shapes.add_picture(img_path1, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.3))
    else:
        # Placeholder panel
        panel = slide_9.shapes.add_shape(1, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.3))
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLOR_LIGHT_BG
        panel.line.color.rgb = COLOR_ROYAL
        
        right_box = slide_9.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.3), Inches(1.0))
        tf_ph = right_box.text_frame
        tf_ph.word_wrap = True
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = "[Polarization Heatmap]"
        p_ph.alignment = PP_ALIGN.CENTER
        p_ph.font.name = FONT_BODY
        p_ph.font.size = Pt(12)
        p_ph.font.color.rgb = COLOR_MUTED

    # Right Table: Community Metrics Table (Bottom)
    table_shape = slide_9.shapes.add_table(6, 5, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.2))
    table = table_shape.table
    table.columns[0].width = Inches(2.3) # Community
    table.columns[1].width = Inches(0.9) # We/Them
    table.columns[2].width = Inches(0.9) # Toxicity
    table.columns[3].width = Inches(0.9) # Pol Index
    table.columns[4].width = Inches(0.8) # Size
    
    headers = ["Community Cluster", "We/Them", "Toxicity", "Polar. Idx", "Size"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_LIGHT_BG
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_ROYAL
        p.alignment = PP_ALIGN.CENTER
        
    table_data = [
        ["Cluster 1 (Progressive Net.)", "1.03", "0.46", "0.473", "353"],
        ["Cluster 3 (Alt-Right Chamber)", "1.07", "0.43", "0.457", "690"],
        ["Cluster 5 (Far-Left Network)", "0.98", "0.43", "0.417", "280"],
        ["Cluster 7 (Local Politics)", "0.68", "0.50", "0.342", "150"],
        ["Cluster 2 (Mainstream Media)", "0.18", "0.48", "0.088", "711"]
    ]
    
    for row_idx, row in enumerate(table_data):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(9.0)
            p.font.color.rgb = COLOR_CHARCOAL
            if col_idx == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.bold = True
            else:
                p.alignment = PP_ALIGN.CENTER

    # Caption below table
    cap_box = slide_9.shapes.add_textbox(Inches(7.0), Inches(6.55), Inches(5.8), Inches(0.4))
    tf_cap = cap_box.text_frame
    tf_cap.word_wrap = True
    p_cap = tf_cap.paragraphs[0]
    p_cap.text = "Figure & Table: Polarization index heatmap (top) and community polarization statistics (bottom)."
    p_cap.alignment = PP_ALIGN.CENTER
    p_cap.font.name = FONT_BODY
    p_cap.font.size = Pt(9.0)
    p_cap.font.italic = True
    p_cap.font.color.rgb = COLOR_MUTED
    
    # ==========================================
    # SLIDE 10: BUT COMPETENCIES & REFLECTION
    # ==========================================
    slide_10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_10)
    add_header(slide_10, "8. BUT Skills & Critical Reflection")
    add_footer(slide_10, 10, total_slides)
    
    left_box = slide_10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_comp_h = tf_left.paragraphs[0]
    p_comp_h.text = "BUT Skill Self-Evaluation"
    p_comp_h.font.name = FONT_TITLE
    p_comp_h.font.size = Pt(18)
    p_comp_h.font.bold = True
    p_comp_h.font.color.rgb = COLOR_ROYAL
    p_comp_h.space_after = Pt(10)
    
    comp_evals = [
        ("Compétence 1 (Développement)", "2.5 -> 4.5 | Structured modules, environment isolation, Git versioning."),
        ("Compétence 2 (Optimisation)", "2.0 -> 4.0 | Configured MPS hardware acceleration for CLIP, WebGL Scattergl for Plotly rendering."),
        ("Compétence 3 (Données)", "3.0 -> 4.5 | Engineered 3-tier schemas under Pandas, GML exports, Louvain network partitioning."),
        ("Compétence 5 (Conduite de Projet)", "2.0 -> 4.0 | Applied Agile timeline iterations, managed weekly supervisor report reviews.")
    ]
    for title, desc in comp_evals:
        p_pt = tf_left.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(11)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(4)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL

    # Right panel - Career & AI Disclosure
    panel_ref = slide_10.shapes.add_shape(1, Inches(6.7), Inches(1.5), Inches(6.133), Inches(5.0))
    panel_ref.fill.solid()
    panel_ref.fill.fore_color.rgb = COLOR_LIGHT_BG
    panel_ref.line.color.rgb = COLOR_ROYAL
    panel_ref.line.width = Pt(1.5)
    
    right_box = slide_10.shapes.add_textbox(Inches(6.85), Inches(1.7), Inches(5.8), Inches(4.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_ref_h = tf_right.paragraphs[0]
    p_ref_h.text = "Career Impact & AI Sources Disclosure"
    p_ref_h.font.name = FONT_TITLE
    p_ref_h.font.size = Pt(18)
    p_ref_h.font.bold = True
    p_ref_h.font.color.rgb = COLOR_ROYAL
    p_ref_h.space_after = Pt(10)
    
    ref_points = [
        ("Career Path Pursuit", "Strengthened my ambition to pursue an engineering school or a research Master's degree in Artificial Intelligence and Big Data."),
        ("AI Utilization Tool", "Antigravity IDE (Gemini 3.5 model) was used for syntax lookup (Streamlit CSS overrides, Plotly layouts) and orthographic checking."),
        ("Intellectual Integrity", "All core code, math indices, community partitions, and texts are 100% human intellectual contributions (<15% assistance).")
    ]
    for title, desc in ref_points:
        p_pt = tf_right.add_paragraph()
        p_pt.text = f"• {title}: "
        p_pt.font.name = FONT_BODY
        p_pt.font.size = Pt(11.5)
        p_pt.font.bold = True
        p_pt.font.color.rgb = COLOR_ROYAL
        p_pt.space_before = Pt(6)
        
        run = p_pt.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10.5)
        run.font.bold = False
        run.font.color.rgb = COLOR_CHARCOAL
        
    # Save presentation
    output_dir = "deliverables"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "RGU_Internship_Defense_Presentation.pptx")
    print(f"Saving presentation to: {output_path}")
    prs.save(output_path)
    print("PowerPoint presentation successfully created.")

if __name__ == "__main__":
    create_presentation()
